""" open() 함수는 close() 명령이 명시되지 않으면 파일이 닫히지 않음. >> 대용량 파일을 처리할 때는 주로 open()을 더 자주 씀.
with open()은 하위 명령줄이 모두 실행되면 파일이 자동으로 닫힘. 그러나 연산능력을 open() 보다 많이 사용. >> 대용량 파일 처리에 부적합.
    ex) with open("파일명", "모드") as f: 로 사용.
------------------------------
파일 입출력 모드 3가지
'r'ead 읽기
    1) f.read() >> 파일 내용 전체를 하나의 str 개체로 반환.
    2) f.readlines() >> 파일의 전체 내용을 각 행을 요소로 갖는 list로 저장.
    3) f.readline() >> 행 단위로 읽어서 반환.
'w'rite 새로 쓰기 >> 모든 내용을 str 형식으로 저장.
'a'dd 추가하기
------------------------------
encoding = 'UTF-8'
위의 명령줄을 파일의 끝에 추가하면 인코딩 설정이 됨.
------------------------------
파일을 읽을 때에는 파일의 위치를 정의해줄 필요가 있음.
1) 절대경로로 정의
    ex - C:\\users\\desktop\\...
2) 같은 위치에 있는 파일로 정의
    ex - 파일 이름만 정의해주면 됨.
------------------------------
"""
def makeinfo():
    f = open("info.txt", "w", encoding = "UTF-8")
    f.write("이름: 윤주영\n나이: 21세\n소속: 숭실대")
    f.close()
    return


def changename():
    with open("info.txt", "r", encoding = "UTF-8") as f:
        contents = f.readlines()
        for i in range(len(contents)):
            if "윤주영" in contents[i]:
                contents[i] = ("이름: 김철수\n") 
    
    with open("info.txt", "w", encoding = "UTF-8") as f:
        for i in range(len(contents)):
            f.write(contents[i])

    return


"""
MS Office 파일 입출력 다루는 법
Excel >> openpyxl
PPT >> Python-pptx
Docx >> Python-docx

Python (Anaconda 아님.)에서 모듈 설치
cmd 상에서 아래 명령어 입력
>>python -m pip install [모듈 이름] = 설치
>>python -m pip install --upgrade [모듈 이름] = 업데이트
주의: python을 cmd 상에서 실행시킨 상태로 진행하지 않고 그냥 cmd를 켜자마자 입력해야 함.
"""
def MakeEXL():
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Hello"
    wb.save("text.xlsx")
    wb.remove(ws)


def MakePPT():
    from pptx import Presentation

    prs = Presentation()

    ts = prs.slide_layouts[0]
    s1 = prs.slides.add_slide(ts)
    title = s1.shapes.title
    st = s1.placeholders[1]
    title.text = "Hello, Python!"
    st.text = "Hello, Python-pptx!"

    prs.save("test.pptx")


if __name__ == "__main__":
    # makeinfo()
    # changename()
    MakeEXL()
    MakePPT()
    pass